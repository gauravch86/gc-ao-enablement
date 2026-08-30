#!/usr/bin/env python3
"""Case Study 2 executive PPTX — sequential 8-slide interview flow (matches CS1 format)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0F, 0x14, 0x19)
SURFACE = RGBColor(0x1A, 0x23, 0x32)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x9C, 0xB3)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
SUCCESS = RGBColor(0x3F, 0xB9, 0x50)
DANGER = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xA3, 0x71, 0xF7)
CYAN = RGBColor(0x38, 0xBD, 0xF8)

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 8


def set_run(run, size=14, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(s, fill)
    try:
        s.adjustments[0] = 0.08
    except Exception:
        pass
    return s


def textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def add_para(tf, text, size=13, bold=False, color=TEXT, align=PP_ALIGN.LEFT, space_after=4):
    if tf.paragraphs[0].text == "" and all(not p.text for p in tf.paragraphs):
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill_shape(bg, BG)
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return slide


def label(slide, text):
    tb = textbox(slide, Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.28))
    add_para(tb.text_frame, text.upper(), size=10, bold=True, color=ACCENT, space_after=0)


def title(slide, text):
    tb = textbox(slide, Inches(0.55), Inches(0.48), Inches(9.7), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, text, size=24, bold=True, color=GOLD, space_after=0)


def footer(slide, page):
    tb = textbox(slide, Inches(0.55), Inches(7.1), Inches(10.2), Inches(0.28))
    add_para(tb.text_frame, "Case Study 2 · Reusable Autonomous Operations offering · Gaurav Chaudhary", size=10, color=MUTED, space_after=0)
    tb2 = textbox(slide, Inches(11.5), Inches(7.1), Inches(1.3), Inches(0.28))
    add_para(tb2.text_frame, f"{page} / {TOTAL}", size=10, color=MUTED, align=PP_ALIGN.RIGHT, space_after=0)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def flow_strip(slide, active):
    steps = ["1 Unlock", "2 Autonomy", "3 Measure", "4 Reuse", "5 Enable", "6 P&L", "7 Ask"]
    y = Inches(6.78)
    width = Inches(1.65)
    gap = Inches(0.12)
    start = Inches(0.55)
    for i, name in enumerate(steps):
        left = start + i * (width + gap)
        on = i == active
        rect(slide, left, y, width, Inches(0.32), SURFACE)
        if on:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, width, Inches(0.05))
            fill_shape(bar, GOLD)
        tb = textbox(slide, left, y + Inches(0.02), width, Inches(0.28))
        add_para(tb.text_frame, name, size=10, bold=on, color=GOLD if on else MUTED, align=PP_ALIGN.CENTER, space_after=0)


def stage_badge(slide, num, name, color=ACCENT):
    rect(slide, Inches(0.55), Inches(1.15), Inches(2.85), Inches(0.42), SURFACE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.15), Inches(0.08), Inches(0.42))
    fill_shape(bar, color)
    tb = textbox(slide, Inches(0.75), Inches(1.2), Inches(2.55), Inches(0.35))
    add_para(tb.text_frame, f"STAGE {num}  ·  {name}", size=11, bold=True, color=color, space_after=0)


def metric_card(slide, left, top, width, height, lbl, val, sub, val_color=GOLD):
    rect(slide, left, top, width, height, SURFACE)
    tb = textbox(slide, left + Inches(0.12), top + Inches(0.12), width - Inches(0.24), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, lbl.upper(), size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER, space_after=2)
    add_para(tf, val, size=22, bold=True, color=val_color, align=PP_ALIGN.CENTER, space_after=2)
    add_para(tf, sub, size=11, color=MUTED, align=PP_ALIGN.CENTER, space_after=0)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ─── 1 Situation ───
    s = blank_slide(prs)
    label(s, "Strategic Product Manager  ·  Case Study 2  ·  10-minute brief  ·  Lived experience")
    title(s, "Reusable Autonomous Operations on customer KPIs / OKRs")

    tb = textbox(s, Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(
        tf,
        "Large telecom operator · multi-vendor · multi-domain · margin and reliability pressure. "
        "Today: reactive operations, high incidents, manual troubleshooting, inconsistent service-level agreements. "
        "Propose a reusable managed-services offering: reactive → predictive → autonomous.",
        size=14, color=MUTED, space_after=0,
    )

    rect(s, Inches(0.55), Inches(2.1), Inches(12.2), Inches(1.2), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(2.25), Inches(11.8), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "NORTH STAR", size=11, bold=True, color=ACCENT, space_after=4)
    add_para(
        tf,
        "End-user experience and Net Promoter Score · service-level / contractual performance-indicator attainment · market trust. "
        "Customer total cost of ownership funds and scales that outcome — it is not the outcome itself.",
        size=13, color=TEXT, space_after=0,
    )

    for i, (h, body, col) in enumerate([
        ("GIVEN", "Ericsson products + third-party portfolio under one Managed Services umbrella", GOLD),
        ("METHOD", "Design → Deploy → Operate → Assure → Improve · pilots before scale", ACCENT),
        ("PROOF", "Telefónica Argentina renewal · Three UK consolidation & autonomy economics", SUCCESS),
    ]):
        left = Inches(0.55) + i * Inches(4.15)
        rect(s, left, Inches(3.55), Inches(3.95), Inches(1.35), SURFACE)
        tb = textbox(s, left + Inches(0.18), Inches(3.7), Inches(3.6), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=12, bold=True, color=col, space_after=6)
        add_para(tf, body, size=13, color=TEXT, space_after=0)

    tb = textbox(s, Inches(0.55), Inches(5.15), Inches(12.2), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "STORYLINE  ·  ONE STAGE PER SLIDE  ·  10 MIN + 5 MIN QUESTIONS", size=11, bold=True, color=GOLD, space_after=8)
    add_para(tf, "1  Economic unlock   →   2  Autonomy path   →   3  Measure from data", size=14, bold=True, color=TEXT, space_after=4)
    add_para(tf, "4  Reuse library   →   5  Enablement cadence   →   6  Five-year profit-and-loss   →   7  Proof & ask", size=14, bold=True, color=TEXT, space_after=6)
    add_para(tf, "Top-right SLIDE BUDGET + gold bar = live countdown in Slideshow (does not auto-advance).", size=12, color=MUTED, space_after=0)
    footer(s, 1)
    set_notes(s, """TIMING: ~1:15

Paint pain in one breath. North star = experience & trust; TCO funds it.
Given: Ericsson + third-party under Managed Services.
Tell them: one stage per slide — lived patterns productized into a reusable offer.

GLOSSARY: KPI · OKR · TCO · SLA · CPI · NPS · 3PP / third-party products.""")

    # ─── 2 Unlock ───
    s = blank_slide(prs)
    label(s, "Storyline step 1 of 7  ·  Economic unlock (enabler, not the north star)")
    title(s, "Vendor consolidation unlocks autonomy — it is not the goal")
    stage_badge(s, "1", "ECONOMIC UNLOCK", GOLD)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.35))
    add_para(tb.text_frame, "Single accountability cuts mean time to restore and customer governance tax", size=13, color=MUTED, space_after=0)

    metric_card(s, Inches(0.55), Inches(1.85), Inches(3.9), Inches(1.55),
                "IT operations budget (assumed)", "$50m", "Multi-vendor run cost today", GOLD)
    metric_card(s, Inches(4.7), Inches(1.85), Inches(3.9), Inches(1.55),
                "After consolidation", "$40m", "~20% total-cost-of-ownership reduction", SUCCESS)
    metric_card(s, Inches(8.85), Inches(1.85), Inches(3.9), Inches(1.55),
                "Customer governance load", "−20%", "Vendor management / service-level overhead", ACCENT)

    for i, (h, lines, col) in enumerate([
        ("PAIN TODAY", ["Finger-pointing across Level-1 / Level-2 / Level-3 / deploy / Service Delivery Manager", "Noisy key performance indicators → more meetings, not resolution", "Cross-domain correlation slow or impossible"], DANGER),
        ("WHAT CONSOLIDATION UNLOCKS", ["One throat to choke from monitor → restore → problem sunset", "Clean baseline for observability and autonomy", "Customer governance shrinks with vendor count"], SUCCESS),
        ("LEADERSHIP CHAIN", ["Level-1 monitor / service desk", "→ Level-2 / Level-3 application", "→ Infrastructure / cloud → change → Service Delivery Manager → problem / sunset"], CYAN),
    ]):
        left = Inches(0.55) + i * Inches(4.15)
        rect(s, left, Inches(3.7), Inches(3.95), Inches(2.55), SURFACE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(3.7), Inches(3.95), Inches(0.06))
        fill_shape(bar, col)
        tb = textbox(s, left + Inches(0.18), Inches(3.9), Inches(3.6), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=12, bold=True, color=col, space_after=8)
        for line in lines:
            add_para(tf, "•  " + line, size=12, color=TEXT, space_after=5)

    flow_strip(s, 0)
    footer(s, 2)
    set_notes(s, """TIMING: ~1:15

SAY: Consolidation is the operating-model unlock — not the north star.
$50m → $40m and −20% governance are economic proof.
Next = autonomy path once single accountability exists.

GLOSSARY: MTTR = mean time to restore · SDM = Service Delivery Manager · TCO = total cost of ownership.""")

    # ─── 3 Autonomy ───
    s = blank_slide(prs)
    label(s, "Storyline step 2 of 7  ·  Outcome path and productized lifecycle")
    title(s, "Reactive → predictive → autonomous — with gates")
    stage_badge(s, "2", "AUTONOMY PATH", ACCENT)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.35))
    add_para(tb.text_frame, "Phased, risk-safe evolution after consolidation — no penalty cliffs", size=13, color=MUTED, space_after=0)

    stages = [
        ("NOW · REACTIVE", "Ticket storms · manual triage · multi-vendor blame", "Gate: consolidate + baseline dump", DANGER),
        ("PHASE A · OBSERVE", "One health view across infrastructure, cloud, application, network", "+ Correlation · single pane of glass", ACCENT),
        ("PHASE B · PREDICT", "Anomaly detection · prevent before end-user impact", "+ Artificial intelligence / machine learning · policy thresholds", PURPLE),
        ("PHASE C · AUTONOMOUS", "Self-heal on known classes · human-in-loop → policy auto", "+ Closed-loop agents · reuse catalog", SUCCESS),
    ]
    for i, (h, body, gate, col) in enumerate(stages):
        left = Inches(0.55) + i * Inches(3.15)
        rect(s, left, Inches(1.85), Inches(3.0), Inches(2.7), SURFACE)
        tb = textbox(s, left + Inches(0.15), Inches(2.0), Inches(2.7), Inches(2.4))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=12, bold=True, color=col, space_after=8)
        add_para(tf, body, size=13, color=TEXT, space_after=10)
        add_para(tf, gate, size=12, color=MUTED, space_after=0)

    rect(s, Inches(0.55), Inches(4.8), Inches(12.2), Inches(1.5), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(4.95), Inches(11.8), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "LIFECYCLE  ·  ONE PRODUCTIZED JOURNEY", size=11, bold=True, color=GOLD, space_after=6)
    add_para(
        tf,
        "Design (intent & policy · key-performance-indicator contract)  →  Deploy (connectors · single pane of glass · shadow mode)  →  "
        "Operate (human-in-the-loop · self-heal)  →  Assure (digital Service Delivery Manager board · credits)  →  "
        "Improve (offender backlog · known-error sunset · library feedback)",
        size=13, color=TEXT, space_after=0,
    )
    flow_strip(s, 1)
    footer(s, 3)
    set_notes(s, """TIMING: ~1:15

Walk four maturity boxes left→right. Lifecycle is one product — AI/orchestration plug into each stage.
GLOSSARY: SPOG = single pane of glass · HITL = human-in-the-loop · Shadow mode · Known error.""")

    # ─── 4 Measure ───
    s = blank_slide(prs)
    label(s, "Storyline step 3 of 7  ·  Data-driven sizing and contractual performance indicators")
    title(s, "Size from incidents. Contract soft-ramp.")
    stage_badge(s, "3", "MEASURE", CYAN)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.35))
    add_para(tb.text_frame, "Do not invent service levels in a workshop — size them from 6–12 months of evidence", size=13, color=MUTED, space_after=0)

    rect(s, Inches(0.55), Inches(1.85), Inches(6.2), Inches(4.4), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(2.05), Inches(5.8), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "DISCOVERY → CONTRACT", size=14, bold=True, color=GOLD, space_after=12)
    for line in [
        "Ingest 6–12 months of incident history",
        "Subject-matter expert + AI enrichment → maturity per component",
        "Top-offender backlog (volume × business impact)",
        "Contractual performance indicators: availability, response/restore, problem aging, change success",
        "Soft-ramp windows — avoid early penalty cliffs",
        "Same scoreboard for roadmap priority and milestone profit-and-loss",
    ]:
        add_para(tf, "•  " + line, size=14, color=TEXT, space_after=8)

    rect(s, Inches(7.0), Inches(1.85), Inches(5.75), Inches(4.4), SURFACE)
    tb = textbox(s, Inches(7.2), Inches(2.05), Inches(5.35), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "ILLUSTRATIVE CPI CREDITS", size=14, bold=True, color=ACCENT, space_after=12)
    for line in [
        "Availability by service class — credits up to ~20%",
        "Priority-0/1 restore ≤4 hours — credit ~10%",
        "Change success ≥98%",
        "Proactive problems ≥20%",
        "Repeat Priority-0–2 ≤10%",
        "Credit % = financial exposure if breached — protects customer return and our margin narrative",
    ]:
        add_para(tf, "•  " + line, size=13, color=TEXT, space_after=8)

    flow_strip(s, 2)
    footer(s, 4)
    set_notes(s, """TIMING: ~1:15

Emphasize soft-ramp and credit-weighted CPIs — commercial gravity, not vanity metrics.
GLOSSARY: CPI · Soft-ramp · Top offender · P0/P1 · Credit.""")

    # ─── 5 Reuse ───
    s = blank_slide(prs)
    label(s, "Storyline step 4 of 7  ·  Architecture for scale across logos")
    title(s, "Productize reuse — next customer inherits the library")
    stage_badge(s, "4", "REUSE LIBRARY", PURPLE)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.35))
    add_para(tb.text_frame, "Why this is a portfolio offering — not a one-account hero project", size=13, color=MUTED, space_after=0)

    items = [
        ("METRICS & RUNBOOK PACKS", "Per domain templates · drag-drop into the next account — not rebuild", ACCENT),
        ("SELF-HEAL CLASSES", "Known failure classes versioned · expand with policy confidence", SUCCESS),
        ("ERICSSON + THIRD-PARTY", "Thin adapters under one Managed Services operating model", GOLD),
        ("TRADE-OFF", "Thin adapters beat rip-and-replace dogma — speed and margin", CYAN),
        ("IMPROVE → GLOBAL", "Offender / known-error feedback compounds the library for every region", PURPLE),
        ("FUNDING", "Account seeds reusable intellectual property; portfolio research & development absorbs long-term platform cost", DANGER),
    ]
    for i, (h, body, col) in enumerate(items):
        row, col_i = divmod(i, 3)
        left = Inches(0.55) + col_i * Inches(4.15)
        top = Inches(1.85) + row * Inches(2.2)
        rect(s, left, top, Inches(3.95), Inches(2.0), SURFACE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.95), Inches(0.06))
        fill_shape(bar, col)
        tb = textbox(s, left + Inches(0.18), top + Inches(0.2), Inches(3.6), Inches(1.65))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=12, bold=True, color=col, space_after=8)
        add_para(tf, body, size=13, color=TEXT, space_after=0)

    flow_strip(s, 3)
    footer(s, 5)
    set_notes(s, """TIMING: ~1:00

This is the productization answer. Adapters for Ericsson and third-party. Improve feeds global library.
GLOSSARY: Adapter · Reuse library · IP = intellectual property.""")

    # ─── 6 Enable ───
    s = blank_slide(prs)
    label(s, "Storyline step 5 of 7  ·  Adoption for customer ops and Ericsson regions")
    title(s, "Enablement cadence — shortened, gated, shared")
    stage_badge(s, "5", "ENABLE", SUCCESS)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.35))
    add_para(tb.text_frame, "Reuse fails without adoption — train, phase, manage resistance, prove success", size=13, color=MUTED, space_after=0)

    # timeline 4 boxes
    phases = [
        ("M0–M3", "Train + baseline · library seed · shadow start"),
        ("M3–M6", "Forward shadow · soft-ramp CPIs · champions live"),
        ("M6–M9", "Reverse shadow · human-in-loop packs · region kit v1"),
        ("M12–M15", "Auto classes expand · Nth-logo pilot · reuse proof"),
    ]
    for i, (m, t) in enumerate(phases):
        left = Inches(0.55) + i * Inches(3.15)
        rect(s, left, Inches(1.85), Inches(3.0), Inches(1.55), SURFACE)
        tb = textbox(s, left + Inches(0.15), Inches(2.0), Inches(2.7), Inches(1.25))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, m, size=14, bold=True, color=ACCENT, space_after=6)
        add_para(tf, t, size=12, color=TEXT, space_after=0)

    for i, (h, lines, col) in enumerate([
        ("CUSTOMER OPERATIONS", ["Role-based training (Level-1 / Level-2–3 / Service Delivery Manager / Problem)", "Shadow → co-pilot → autonomy coverage expand", "Reframe “jobs cut” as A-team / problem ownership upskill", "Success: soft-ramp CPIs met · ticket quality ≥95%"], ACCENT),
        ("ERICSSON REGIONS", ["Train-the-trainer on maturity framework and Deal Desk tiers", "Configure-from-library first; custom via exception board", "Kit: one-pager · demo · playbooks · pricing cards", "Success: time-to-first-value · % library reuse · margin"], PURPLE),
        ("PILOTS BEFORE SCALE", ["Validate total-cost trajectory and soft-ramp credibility", "Pass-gates before regional rollout", "M9–M12 = human-in-loop stabilize before auto expand", "Protects customer trust and our profit-and-loss"], GOLD),
    ]):
        left = Inches(0.55) + i * Inches(4.15)
        rect(s, left, Inches(3.65), Inches(3.95), Inches(2.6), SURFACE)
        tb = textbox(s, left + Inches(0.18), Inches(3.8), Inches(3.6), Inches(2.3))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=12, bold=True, color=col, space_after=8)
        for line in lines:
            add_para(tf, "•  " + line, size=12, color=TEXT, space_after=4)

    flow_strip(s, 4)
    footer(s, 6)
    set_notes(s, """TIMING: ~1:15

Hit the shortened map: Reverse shadow M6–M9 · Auto classes M12–M15.
M9–M12 intentional stabilize window.
Two audiences: customer ops + Ericsson regions.
GLOSSARY: Forward/Reverse shadow · HITL · Pass-gate · Nth logo.""")

    # ─── 7 P&L ───
    s = blank_slide(prs)
    label(s, "Storyline step 6 of 7  ·  Five-year managed-services profit-and-loss")
    title(s, "Full-time-equivalent ramp · site-mix · tools · penalty lot")
    stage_badge(s, "6", "P&L CONSTRUCT", GOLD)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(
        tf,
        "Baseline multi-vendor headcount = B. Target X = B − 30%. Customer commit ~20% run-cost cut; we size at 30% → ~10-point structural earnings gap.",
        size=13, color=MUTED, space_after=0,
    )

    years = [("Y1", "X+20%", "Knowledge-transfer peak"), ("Y2", "X+10%", "Human-in-loop expand"),
             ("Y3", "X", "Steady-state"), ("Y4", "X−10%", "Autonomy harvest"), ("Y5", "X−20%", "Closed-loop scale")]
    for i, (y, v, r) in enumerate(years):
        left = Inches(0.55) + i * Inches(2.5)
        rect(s, left, Inches(1.9), Inches(2.35), Inches(1.55), SURFACE)
        tb = textbox(s, left + Inches(0.1), Inches(2.05), Inches(2.15), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, y, size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, space_after=2)
        add_para(tf, v, size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER, space_after=2)
        add_para(tf, r, size=11, color=MUTED, align=PP_ALIGN.CENTER, space_after=0)

    levers = [
        ("① Labour gap", "+10%", "Size −30% vs commit −20%", SUCCESS),
        ("② Site-mix", "+10%", "Onsite leads · offshore factory", SUCCESS),
        ("③ Tools seed", "−5%", "Reusable IP start cost", DANGER),
        ("④ License recover", "+5%", "Customer redirects tooling budget", SUCCESS),
        ("⑤ Penalty lot", "+5% if met", "Parked envelope → earn-back", SUCCESS),
    ]
    for i, (h, v, sub, col) in enumerate(levers):
        left = Inches(0.55) + i * Inches(2.5)
        rect(s, left, Inches(3.7), Inches(2.35), Inches(2.55), SURFACE)
        tb = textbox(s, left + Inches(0.1), Inches(3.9), Inches(2.15), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER, space_after=6)
        add_para(tf, v, size=18, bold=True, color=col, align=PP_ALIGN.CENTER, space_after=6)
        add_para(tf, sub, size=12, color=MUTED, align=PP_ALIGN.CENTER, space_after=0)

    flow_strip(s, 5)
    footer(s, 7)
    set_notes(s, """TIMING: ~1:30

Deal Desk lean-in slide. Plan 10–15% EBIT with envelope; ~20% upside earned if CPIs hold.
GLOSSARY: FTE · B · X · EBIT · Site-mix · Penalty lot · Earn-back.""")

    # ─── 8 Ask ───
    s = blank_slide(prs)
    label(s, "Storyline step 7 of 7  ·  Lived proof and recommendation")
    title(s, "Lived proof. One recommendation.")
    stage_badge(s, "7", "PROOF & ASK", SUCCESS)

    rect(s, Inches(0.55), Inches(1.8), Inches(6.0), Inches(2.35), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(1.95), Inches(5.6), Inches(2.05))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "TELEFÓNICA ARGENTINA — ~$100m RENEWAL", size=12, bold=True, color=GOLD, space_after=8)
    add_para(tf, "3,000+ backlog · ~500/day → offender taxonomy", size=13, color=TEXT, space_after=4)
    add_para(tf, "Executive key-performance-indicator transparency · 24×7 scale", size=13, color=TEXT, space_after=4)
    add_para(tf, "Accountability restored → renewal won on trust", size=13, color=TEXT, space_after=0)

    rect(s, Inches(6.8), Inches(1.8), Inches(5.95), Inches(2.35), SURFACE)
    tb = textbox(s, Inches(7.0), Inches(1.95), Inches(5.55), Inches(2.05))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "THREE UK — ~$125m VENDOR CONSOLIDATION", size=12, bold=True, color=ACCENT, space_after=8)
    add_para(tf, "Four vendors · 400+ components · ~75 contractual performance indicators", size=13, color=TEXT, space_after=4)
    add_para(tf, "Single pane of glass + self-heal + generative-AI runbooks", size=13, color=TEXT, space_after=4)
    add_para(tf, "474→336 full-time-equivalent economics · lowest Priority-1s", size=13, color=TEXT, space_after=0)

    rect(s, Inches(0.55), Inches(4.4), Inches(12.2), Inches(1.85), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(4.55), Inches(11.8), Inches(1.55))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "RECOMMENDATION", size=11, bold=True, color=SUCCESS, space_after=6)
    add_para(
        tf,
        "Launch a reusable Autonomous Operations managed-services offering (Ericsson + third-party under Managed Services) "
        "that sells autonomy maturity + unified observability + outcome key performance indicators as one productized journey — "
        "with enablement for customer operations and Ericsson regions, and pilots that validate total cost of ownership, "
        "soft-ramp, and library reuse before scaling. Enablement: Reverse shadow M6–M9 · Auto classes M12–M15.",
        size=13, color=TEXT, space_after=0,
    )

    tb = textbox(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.3))
    add_para(
        tb.text_frame,
        "QUESTIONS  ·  5 MINUTES     ·     Deep dive on any stage: HTML case study",
        size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER, space_after=0,
    )
    flow_strip(s, 6)
    footer(s, 8)
    set_notes(s, """TIMING: ~1:00 then Q&A

Two proof anchors then recommendation sentence. Invite depth on CPIs, five-year economics, or enablement.

LIKELY Qs:
• Why consolidation first? Unlock — not north star.
• Penalty shock? Soft-ramp + data-sized CPIs + penalty lot.
• Reusable how? Library + thin adapters.
• Pilot fails a gate? Do not scale.""")

    out = "/workspace/case-study-2-ao-offering-executive.pptx"
    prs.save(out)
    from inject_slide_timers import apply_timers
    # 1:15×4 + 1:00 + 1:15 + 1:30 + 1:15 = 10:00
    apply_timers(out, [75, 75, 75, 75, 60, 75, 90, 75], talk_total_sec=600)
    print(f"Wrote {out} ({TOTAL} slides) with live slide timers")
    return out


if __name__ == "__main__":
    build()
