#!/usr/bin/env python3
"""Case Study 1 executive PPTX — sequential interview flow (8 slides), no on-slide timers."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0F, 0x14, 0x19)
SURFACE = RGBColor(0x1A, 0x23, 0x32)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x9C, 0xB3)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
SUCCESS = RGBColor(0x3F, 0xB9, 0x50)
DANGER = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xA3, 0x71, 0xF7)

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 8


def set_run(run, size=14, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(s, fill)
    try:
        s.adjustments[0] = 0.08
    except Exception:
        pass
    return s


def textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def add_para(tf, text, size=13, bold=False, color=TEXT, align=PP_ALIGN.LEFT, space_after=4):
    if tf.paragraphs[0].text == "" and all(not p.text for p in tf.paragraphs):
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill_shape(bg, BG)
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return slide


def label(slide, text):
    tb = textbox(slide, Inches(0.55), Inches(0.22), Inches(12.2), Inches(0.28))
    add_para(tb.text_frame, text.upper(), size=10, bold=True, color=ACCENT, space_after=0)


def title(slide, text, top=Inches(0.48)):
    tb = textbox(slide, Inches(0.55), top, Inches(12.2), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, text, size=24, bold=True, color=GOLD, space_after=0)


def footer(slide, page):
    tb = textbox(slide, Inches(0.55), Inches(7.1), Inches(10), Inches(0.28))
    add_para(
        tb.text_frame,
        "Case Study 1 · Reposition underperforming offering · Gaurav Chaudhary",
        size=10,
        color=MUTED,
        space_after=0,
    )
    tb2 = textbox(slide, Inches(11.5), Inches(7.1), Inches(1.3), Inches(0.28))
    add_para(tb2.text_frame, f"{page} / {TOTAL}", size=10, color=MUTED, align=PP_ALIGN.RIGHT, space_after=0)


def metric_card(slide, left, top, width, height, lbl, val, sub, val_color=GOLD):
    rect(slide, left, top, width, height, SURFACE)
    tb = textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, lbl.upper(), size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER, space_after=2)
    add_para(tf, val, size=22, bold=True, color=val_color, align=PP_ALIGN.CENTER, space_after=2)
    add_para(tf, sub, size=11, color=MUTED, align=PP_ALIGN.CENTER, space_after=0)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def flow_strip(slide, active):
    steps = [
        "1 Diagnose",
        "2 Findings",
        "3 Decide",
        "4 Package",
        "5 Guardrails",
        "6 Align",
        "7 Ask",
    ]
    y = Inches(6.78)
    width = Inches(1.65)
    gap = Inches(0.12)
    start = Inches(0.55)
    for i, name in enumerate(steps):
        left = start + i * (width + gap)
        on = i == active
        rect(slide, left, y, width, Inches(0.32), SURFACE)
        if on:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y, width, Inches(0.05))
            fill_shape(bar, GOLD)
        tb = textbox(slide, left, y + Inches(0.02), width, Inches(0.28))
        add_para(
            tb.text_frame,
            name,
            size=10,
            bold=on,
            color=GOLD if on else MUTED,
            align=PP_ALIGN.CENTER,
            space_after=0,
        )


def stage_badge(slide, num, name, color=ACCENT):
    rect(slide, Inches(0.55), Inches(1.15), Inches(2.85), Inches(0.42), SURFACE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.15), Inches(0.08), Inches(0.42))
    fill_shape(bar, color)
    tb = textbox(slide, Inches(0.75), Inches(1.2), Inches(2.55), Inches(0.35))
    add_para(tb.text_frame, f"STAGE {num}  ·  {name}", size=11, bold=True, color=color, space_after=0)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ─── SLIDE 1 · Situation ───
    s = blank_slide(prs)
    label(s, "Strategic Product Manager  ·  Case Study 1  ·  10-minute commercial brief")
    title(s, "Reposition an underperforming portfolio offering")

    tb = textbox(s, Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(
        tf,
        "Application Managed Services: strong technical delivery, but commercially stuck — "
        "flat net sales, declining gross margin, Sales cannot articulate value, inconsistent regional pricing.",
        size=14,
        color=MUTED,
        space_after=0,
    )

    metric_card(
        s,
        Inches(0.55),
        Inches(1.95),
        Inches(3.9),
        Inches(1.45),
        "Net sales (2 years)",
        "Flat",
        "Weak conversion · little attach",
        MUTED,
    )
    metric_card(
        s,
        Inches(4.7),
        Inches(1.95),
        Inches(3.9),
        Inches(1.45),
        "Gross margin (illustrative)",
        "28% → 22–24%",
        "−4 to −6 percentage points over 2 years",
        DANGER,
    )
    metric_card(
        s,
        Inches(8.85),
        Inches(1.95),
        Inches(3.9),
        Inches(1.45),
        "Delivery satisfaction",
        "Strong",
        "Execution is not the failure mode",
        SUCCESS,
    )

    rect(s, Inches(0.55), Inches(3.65), Inches(12.2), Inches(1.15), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(3.8), Inches(11.8), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "EXECUTIVE FRAMING", size=11, bold=True, color=ACCENT, space_after=4)
    add_para(
        tf,
        "Commercial / go-to-market problem — not a delivery turnaround. Sold as headcount; buyers purchase outcomes, "
        "cost predictability, and risk reduction. Plan must restore net sales and gross margin together — "
        "discounts only when traded for scope, term, or attach. Leadership asks: evolve, reposition, or sunset.",
        size=13,
        color=TEXT,
        space_after=0,
    )

    tb = textbox(s, Inches(0.55), Inches(5.05), Inches(12.2), Inches(1.55))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "STORYLINE  ·  ONE STAGE PER SLIDE  ·  10 MIN + 5 MIN QUESTIONS", size=11, bold=True, color=GOLD, space_after=8)
    add_para(tf, "1  Five-lens diagnosis", size=14, bold=True, color=TEXT, space_after=3)
    add_para(tf, "2  Illustrative findings (what the data would show)", size=14, bold=True, color=TEXT, space_after=3)
    add_para(tf, "3  Decision framework → Reposition", size=14, bold=True, color=TEXT, space_after=3)
    add_para(
        tf,
        "4  Package & price   ·   5  Guardrails & invest   ·   6  Align & measure   ·   7  Ask",
        size=14,
        bold=True,
        color=TEXT,
        space_after=0,
    )
    footer(s, 1)
    set_notes(
        s,
        """TIMING: ~1:15

Paradox first: delivery satisfaction strong, sales flat, margin illustratively 28% → 22–24%.
Commercial problem, not delivery. Storyline is sequential — one stage per slide.
In framing: we will fix packaging/pricing so net sales and gross margin recover together;
discounts are a trade for scope/term/attach — not free bleed.""",
    )

    # ─── SLIDE 2 · Diagnosis ───
    s = blank_slide(prs)
    label(s, "Storyline step 1 of 7  ·  Evidence before any portfolio call")
    title(s, "Five-lens diagnosis — what we would gather")
    stage_badge(s, "1", "FIVE-LENS DIAGNOSIS", ACCENT)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.35))
    add_para(
        tb.text_frame,
        "Four-week commercial diagnostic · parallel workstreams · one steering forum",
        size=13,
        color=MUTED,
        space_after=0,
    )

    lenses = [
        (
            "A  WIN / LOSS",
            "Top 10 accounts.\nTag losses: price-only\nvs value unclear.\nTrigger: value unclear\n>30% → positioning broken.",
        ),
        (
            "B  MARGIN",
            "GM waterfall: discount\n· transition · unpaid scope.\nAsk: was discount traded\nfor scope / term / attach?\nBleed without trade → fix.",
        ),
        (
            "C  COMPETITIVE",
            "Staff augmentation vs\noutcome-packaged peers.\nTrigger: buyers ask for\ntiers → reposition.",
        ),
        (
            "D  CUSTOMER VOICE",
            "High satisfaction +\nlow willingness to pay =\nstory gap, not\nexecution gap.",
        ),
        (
            "E  PORTFOLIO FIT",
            "18-month path to\ntopline + margin?\nPull-through?\nSunset only if both fail.",
        ),
    ]
    for i, (h, body) in enumerate(lenses):
        left = Inches(0.55) + i * Inches(2.5)
        rect(s, left, Inches(1.75), Inches(2.35), Inches(4.55), SURFACE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.75), Inches(2.35), Inches(0.08))
        fill_shape(bar, ACCENT)
        tb = textbox(s, left + Inches(0.12), Inches(2.0), Inches(2.1), Inches(4.15))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=12, bold=True, color=ACCENT, space_after=8)
        for line in body.split("\n"):
            add_para(tf, line, size=12, color=TEXT, space_after=3)

    flow_strip(s, 0)
    footer(s, 2)
    set_notes(
        s,
        """TIMING: ~1:15

Walk A→E. Win/loss = top 10 accounts (not random 15–20 deals).
On Margin: emphasize whether discounts bought something (scope, term, attach) or were free bleed.
Portfolio fit now includes path to BOTH topline and margin.""",
    )

    # ─── SLIDE 3 · Findings ───
    s = blank_slide(prs)
    label(s, "Storyline step 2 of 7  ·  Illustrative synthesis from the diagnostic")
    title(s, "Root causes — commercial, not delivery")
    stage_badge(s, "2", "ILLUSTRATIVE FINDINGS", DANGER)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.35))
    add_para(
        tb.text_frame,
        "Three failures that kill new sales and erode EBIT at the same time",
        size=13,
        color=MUTED,
        space_after=0,
    )

    causes = [
        (
            "1  WRONG SELLABLE UNIT",
            "Sold as L1/L2/L3 staffing towers.",
            "Procurement benches rate-per-person — new logos stall.",
            "Hits topline: flat net sales · weak attach.",
        ),
        (
            "2  DISCOUNT WITHOUT TRADE",
            "Regions discount to hit quota — no floor, no scope trade.",
            "Should be: discount sustained revenue only when traded for more scope, term, or investment.",
            "Hits margin: GM 28% → 22–24% while NPS holds.",
        ),
        (
            "3  VALUE STORY GAP",
            "Buyers want outcomes and predictability.",
            "Sales still pitches headcount replacement.",
            "Hits both: cannot win new revenue or defend price.",
        ),
    ]
    for i, (h, a, b, c) in enumerate(causes):
        left = Inches(0.55) + i * Inches(4.15)
        rect(s, left, Inches(1.75), Inches(3.95), Inches(4.55), SURFACE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.75), Inches(3.95), Inches(0.08))
        fill_shape(bar, DANGER)
        tb = textbox(s, left + Inches(0.2), Inches(2.0), Inches(3.55), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=13, bold=True, color=DANGER, space_after=12)
        add_para(tf, a, size=14, color=TEXT, space_after=8)
        add_para(tf, b, size=13, color=TEXT, space_after=8)
        add_para(tf, c, size=13, bold=True, color=GOLD, space_after=0)

    flow_strip(s, 1)
    footer(s, 3)
    set_notes(
        s,
        """TIMING: ~1:20

Make discount point land hard: we do NOT ban discounts — we ban discounts that buy nothing.
Proper use: sustain topline while negotiating additional scope, multi-year term, attach packs, or customer-funded enablement.
All three causes are commercial; delivery is fine.""",
    )

    # ─── SLIDE 4 · Decision ───
    s = blank_slide(prs)
    label(s, "Storyline step 3 of 7  ·  Portfolio call with explicit trade-offs")
    title(s, "Decision framework → Reposition")
    stage_badge(s, "3", "DECISION FRAMEWORK", SUCCESS)

    tb = textbox(s, Inches(3.55), Inches(1.2), Inches(9.2), Inches(0.35))
    add_para(
        tb.text_frame,
        "Change what Sales sells and how Deal Desk prices — tooling alone will not fix conversion or margin",
        size=13,
        color=MUTED,
        space_after=0,
    )

    decisions = [
        (
            "SUNSET",
            MUTED,
            False,
            [
                "Rejects delivery DNA and account footprint",
                "Competitors inherit the revenue base",
                "Only if no 18-month path to sales + margin AND no pull-through",
                "Not indicated here",
            ],
        ),
        (
            "EVOLVE ONLY",
            MUTED,
            False,
            [
                "More tooling — same sellable unit",
                "Does not fix Sales narrative or new-logo conversion",
                "Does not stop discount-without-trade",
                "Necessary later; insufficient alone",
            ],
        ),
        (
            "REPOSITION  ✓",
            SUCCESS,
            True,
            [
                "Outcome-tiered Application Operations",
                "Stabilize → Optimize → Transform",
                "New value prop + pricing + Deal Desk rules",
                "Trade-off: near-term deal size may dip; attach + renewals + cleaner GM recover",
            ],
        ),
    ]
    for i, (h, color, selected, lines) in enumerate(decisions):
        left = Inches(0.55) + i * Inches(4.15)
        rect(s, left, Inches(1.75), Inches(3.95), Inches(4.55), SURFACE)
        if selected:
            border = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.75), Inches(3.95), Inches(4.55))
            border.fill.background()
            border.line.color.rgb = SUCCESS
            border.line.width = Pt(2.5)
            try:
                border.adjustments[0] = 0.08
            except Exception:
                pass
        tb = textbox(s, left + Inches(0.2), Inches(2.0), Inches(3.55), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, h, size=16, bold=True, color=color, space_after=12)
        for line in lines:
            add_para(tf, "•  " + line, size=13, color=TEXT if selected else MUTED, space_after=7)

    flow_strip(s, 2)
    footer(s, 4)
    set_notes(
        s,
        """TIMING: ~1:10

Recommendation: REPOSITION. Explicitly dual outcome — new sales growth AND margin recovery.
Evolve-only fails both levers. Sunset destroys the footprint SPM would grow.""",
    )

    # ─── SLIDE 5 · Package & value prop ───
    s = blank_slide(prs)
    label(s, "Storyline step 4 of 7  ·  Revised packaging, pricing & value proposition")
    title(s, "What Sales sells — tiers, outcomes, clearer price")
    stage_badge(s, "4", "PACKAGE & PRICE", GOLD)

    tiers = [
        (
            "STABILIZE",
            GOLD,
            "Land — restore control",
            [
                "Backlog · service-level baseline",
                "Single accountability",
                "Platform fee + criticality class",
                "Floor GM ≥ 22% · entry for new logos",
            ],
        ),
        (
            "OPTIMIZE",
            ACCENT,
            "Expand — prove efficiency",
            [
                "+ Monitoring / automation packs",
                "Outcome-band vs service levels",
                "Optional gain-share on savings",
                "Floor GM ≥ 26% · attach engine",
            ],
        ),
        (
            "TRANSFORM",
            SUCCESS,
            "Deepen — operating uplift",
            [
                "+ Continuous improvement",
                "Milestone elements (Finance OK)",
                "Executive scorecard",
                "Floor GM ≥ 28% · premium renewals",
            ],
        ),
    ]
    for i, (name, color, tagline, lines) in enumerate(tiers):
        left = Inches(0.55) + i * Inches(4.15)
        rect(s, left, Inches(1.7), Inches(3.95), Inches(3.35), SURFACE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.7), Inches(3.95), Inches(0.08))
        fill_shape(bar, color)
        tb = textbox(s, left + Inches(0.2), Inches(1.95), Inches(3.55), Inches(2.95))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, name, size=16, bold=True, color=color, space_after=3)
        add_para(tf, tagline, size=12, bold=True, color=GOLD, space_after=8)
        for line in lines:
            add_para(tf, "•  " + line, size=12, color=TEXT, space_after=4)

    rect(s, Inches(0.55), Inches(5.2), Inches(12.2), Inches(1.3), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(5.32), Inches(11.8), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "UPDATED VALUE PROPOSITION  ·  ONE LINE FOR SALES", size=11, bold=True, color=ACCENT, space_after=4)
    add_para(
        tf,
        "“Ericsson Application Operations — outcome-tiered managed applications that restore control, "
        "improve total cost of ownership where contracted, and give Sales a globally consistent story "
        "with Deal Desk pricing discipline that grows revenue without giving margin away.”",
        size=13,
        color=TEXT,
        space_after=0,
    )
    flow_strip(s, 3)
    footer(s, 5)
    set_notes(
        s,
        """TIMING: ~1:25

This slide answers: revised packaging + pricing + updated value proposition.
Ladder: Stabilize lands new logos → Optimize expands → Transform deepens.
Each tier has a GM floor — topline growth with margin discipline.""",
    )

    # ─── SLIDE 6 · Guardrails + invest ───
    s = blank_slide(prs)
    label(s, "Storyline step 5 of 7  ·  Commercial guardrails + investment ROI sequence")
    title(s, "Deal Desk rules + what we fund (and stop)")
    stage_badge(s, "5", "GUARDRAILS & INVEST", PURPLE)

    rect(s, Inches(0.55), Inches(1.7), Inches(6.3), Inches(4.6), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(1.85), Inches(5.9), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "COMMERCIAL GUARDRAILS  ·  DEAL DESK", size=12, bold=True, color=GOLD, space_after=8)
    for name, detail in [
        ("Floor gross margin", "Stabilize ≥22% · Optimize ≥26% · Transform ≥28%"),
        ("Discount = trade, not gift", "≤12% off list only if traded for scope, term ≥24 mo, or attach pack"),
        ("Sustain topline rule", "If discount needed to keep/win revenue → must increase contracted value or commitment"),
        ("Transition pricing", "Knowledge-transfer mandatory · min 8% Year-1 contract value"),
        ("Scope change", "Unpaid work cap 2% of annual contract value"),
        ("Regional variance", "Like-for-like within ±15% of global rate card"),
        ("Credits / penalties", "Soft-ramp M1–6 · Finance revenue-recognition OK"),
    ]:
        add_para(tf, name, size=11, bold=True, color=ACCENT, space_after=1)
        add_para(tf, detail, size=11, color=TEXT, space_after=5)

    rect(s, Inches(7.05), Inches(1.7), Inches(5.7), Inches(2.55), SURFACE)
    tb = textbox(s, Inches(7.25), Inches(1.85), Inches(5.3), Inches(2.25))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "FUND FIRST  ·  MAX ROI", size=12, bold=True, color=SUCCESS, space_after=6)
    for line in [
        "Q1: Rate card · tier kit · Deal Desk playbook",
        "Q2: Service-level catalogue · 2 regional pilots",
        "Q3–4: Attach packs · renewal step-up motions",
        "Year 2: Retire headcount-only contracts",
    ]:
        add_para(tf, "▸  " + line, size=12, color=TEXT, space_after=5)

    rect(s, Inches(7.05), Inches(4.4), Inches(5.7), Inches(1.9), SURFACE)
    tb = textbox(s, Inches(7.25), Inches(4.55), Inches(5.3), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "DEPRIORITIZE / STOP", size=12, bold=True, color=DANGER, space_after=6)
    for line in [
        "Custom tooling per account before global kit",
        "Below-floor proofs of concept",
        "Regional bespoke rate cards",
        "Rebrand without Sales / Deal Desk kit",
    ]:
        add_para(tf, "×  " + line, size=12, color=TEXT, space_after=3)

    flow_strip(s, 4)
    footer(s, 6)
    set_notes(
        s,
        """TIMING: ~1:20

This slide answers: commercial guardrails/pricing guidelines AND investment prioritization.
Discount rule to verbalize: discount protects topline only when negotiated for more scope, longer term, or attach — increasing total contracted value / investment, not giving margin away.
Fund commercial assets before tooling. Explicit stop list = ROI discipline.""",
    )

    # ─── SLIDE 7 · Align + metrics ───
    s = blank_slide(prs)
    label(s, "Storyline step 6 of 7  ·  Cross-functional plan + success metrics + feedback loop")
    title(s, "One program — Sales, Marketing, Finance, regions")
    stage_badge(s, "6", "ALIGN & MEASURE", ACCENT)

    funcs = [
        ("SALES / DEAL DESK", ACCENT, "≥80% bids on tier + rate card by Q2", "Discovery kit · ROI calc · discount-trade checklist"),
        ("MARKETING", PURPLE, "One narrative: App Ops by tier", "One-pagers · proof · retire staffing language"),
        ("FINANCE", GOLD, "Rev-rec · margin floors · EBIT view", "Platform vs milestone memo · bid sign-off"),
        ("REGIONS / ENABLE", SUCCESS, "Pilots hit catalogue · renewals step up", "KT playbook · regional train-the-trainer"),
    ]
    for i, (name, color, outcome, arts) in enumerate(funcs):
        left = Inches(0.55) + i * Inches(3.15)
        rect(s, left, Inches(1.7), Inches(3.0), Inches(2.35), SURFACE)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.7), Inches(0.08), Inches(2.35))
        fill_shape(bar, color)
        tb = textbox(s, left + Inches(0.18), Inches(1.85), Inches(2.65), Inches(2.05))
        tf = tb.text_frame
        tf.word_wrap = True
        add_para(tf, name, size=11, bold=True, color=color, space_after=5)
        add_para(tf, outcome, size=12, bold=True, color=TEXT, space_after=5)
        add_para(tf, arts, size=11, color=MUTED, space_after=0)

    rect(s, Inches(0.55), Inches(4.25), Inches(12.2), Inches(2.15), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(4.4), Inches(11.8), Inches(1.85))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "12-MONTH SCOREBOARD  ·  Quarterly Pricing Council", size=12, bold=True, color=GOLD, space_after=6)
    add_para(
        tf,
        "Net sales Flat → +8–12%   ·   Gross margin recover 3–4 pts (e.g. toward 25–27%)   ·   Tier adoption ≥80%",
        size=13,
        color=TEXT,
        space_after=4,
    )
    add_para(
        tf,
        "Price variance <15%   ·   Wins on value/TCO ≥40%   ·   Renewal step-up ≥25%   ·   ≥70% discounts have a documented scope/term trade",
        size=13,
        color=TEXT,
        space_after=6,
    )
    add_para(
        tf,
        "Feedback loop: monthly Deal Desk overrides → quarterly Pricing Council (±3% rate card, pack prices) → annual sunset / reposition re-test",
        size=12,
        color=MUTED,
        space_after=0,
    )

    flow_strip(s, 5)
    footer(s, 7)
    set_notes(
        s,
        """TIMING: ~1:10

Answers: cross-functional alignment + success metrics + continuous feedback loop.
Scoreboard tracks net sales and gross margin recovery plus commercial hygiene (tier adoption, discount trades).
Pricing Council = how pricing/commercial models stay alive from market performance.""",
    )

    # ─── SLIDE 8 · Ask ───
    s = blank_slide(prs)
    label(s, "Storyline step 7 of 7  ·  Close and hand to questions")
    title(s, "Reposition. Land with tiers. Expand when value is proven.")
    stage_badge(s, "7", "ASK", SUCCESS)

    rect(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(1.5), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(1.85), Inches(11.8), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "THE CALL", size=11, bold=True, color=ACCENT, space_after=5)
    add_para(
        tf,
        "Reposition Application Managed Services as Outcome-Tiered Application Operations "
        "(Stabilize → Optimize → Transform). Do not sunset. Do not evolve tooling without fixing the sellable unit. "
        "Use discounts to sustain net sales only when traded for scope, term, or attach. "
        "Fund Deal Desk guardrails in Q1; pilot in two regions; govern with a quarterly Pricing Council.",
        size=13,
        color=TEXT,
        space_after=0,
    )

    rect(s, Inches(0.55), Inches(3.4), Inches(6.0), Inches(2.85), SURFACE)
    tb = textbox(s, Inches(0.75), Inches(3.55), Inches(5.6), Inches(2.55))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "THOUGHT PROCESS THEY SHOULD HEAR", size=11, bold=True, color=GOLD, space_after=6)
    for line in [
        "1 Diagnose with five lenses",
        "2 Synthesize commercial root causes",
        "3 Decide with trade-offs → Reposition",
        "4 Package / guardrail / enable / measure",
    ]:
        add_para(tf, line, size=13, color=TEXT, space_after=5)

    rect(s, Inches(6.8), Inches(3.4), Inches(5.95), Inches(2.85), SURFACE)
    tb = textbox(s, Inches(7.0), Inches(3.55), Inches(5.55), Inches(2.55))
    tf = tb.text_frame
    tf.word_wrap = True
    add_para(tf, "ASK OF LEADERSHIP", size=11, bold=True, color=SUCCESS, space_after=6)
    for line in [
        "Approve reposition + global rate card",
        "Stand up Deal Desk trade rules this quarter",
        "Authorize two regional pilots",
        "Name Product owner + monthly steering",
    ]:
        add_para(tf, "▸  " + line, size=13, color=TEXT, space_after=5)

    flow_strip(s, 6)
    footer(s, 8)
    set_notes(
        s,
        """TIMING: ~1:00 then Q&A

Restate call. Point at the four-step thought process — that is the interview signal.
Then four asks. Hand to questions.

Q&A: discount-as-trade · sunset · deal size dip · Finance rev-rec · optional multi-vendor pack.""",
    )

    out = "/workspace/case-study-1-reposition-executive.pptx"
    prs.save(out)
    print(f"Wrote {out} ({TOTAL} slides) — no on-slide timers")
    return out


if __name__ == "__main__":
    build()
