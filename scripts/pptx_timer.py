#!/usr/bin/env python3
"""On-slide presentation timer HUD for interview decks (python-pptx).

Adds a visible slide-budget badge plus a countdown bar that starts when the
slide appears in Slideshow mode. Does NOT auto-advance — presenter stays in control.
"""

from __future__ import annotations

from lxml import etree
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SURFACE = RGBColor(0x1A, 0x23, 0x32)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x9C, 0xB3)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
DANGER = RGBColor(0xF8, 0x71, 0x71)
TRACK = RGBColor(0x24, 0x30, 0x44)


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_run(run, size=12, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _para(tf, text, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, space_after=0):
    if tf.paragraphs[0].text == "" and all(not p.text for p in tf.paragraphs):
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return p


def fmt_mmss(seconds: int) -> str:
    seconds = max(0, int(seconds))
    return f"{seconds // 60}:{seconds % 60:02d}"


def add_countdown_wipe(slide, bar_shape, nudge_shape, dur_ms: int) -> None:
    """Countdown wipe on bar, then reveal MOVE ON nudge after budget (Slideshow)."""
    bar_id = str(bar_shape.shape_id)
    nudge_id = str(nudge_shape.shape_id)
    existing = slide._element.find(qn("p:timing"))
    if existing is not None:
        slide._element.remove(existing)

    # Hide nudge immediately, wipe bar over dur_ms, then fade nudge in (still no auto-advance).
    xml = f"""
    <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="4" fill="hold">
                              <p:stCondLst>
                                <p:cond delay="0"/>
                              </p:stCondLst>
                              <p:childTnLst>
                                <p:set>
                                  <p:cBhvr>
                                    <p:cTn id="5" dur="1" fill="hold"/>
                                    <p:tgtEl>
                                      <p:spTgt spid="{nudge_id}"/>
                                    </p:tgtEl>
                                    <p:attrNameLst>
                                      <p:attrName>style.visibility</p:attrName>
                                    </p:attrNameLst>
                                  </p:cBhvr>
                                  <p:to>
                                    <p:strVal val="hidden"/>
                                  </p:to>
                                </p:set>
                                <p:animEffect transition="out" filter="wipe(left)">
                                  <p:cBhvr>
                                    <p:cTn id="6" dur="{dur_ms}" fill="hold"/>
                                    <p:tgtEl>
                                      <p:spTgt spid="{bar_id}"/>
                                    </p:tgtEl>
                                  </p:cBhvr>
                                </p:animEffect>
                                <p:animEffect transition="in" filter="fade">
                                  <p:cBhvr>
                                    <p:cTn id="7" dur="400" fill="hold">
                                      <p:stCondLst>
                                        <p:cond delay="{dur_ms}"/>
                                      </p:stCondLst>
                                    </p:cTn>
                                    <p:tgtEl>
                                      <p:spTgt spid="{nudge_id}"/>
                                    </p:tgtEl>
                                  </p:cBhvr>
                                </p:animEffect>
                                <p:set>
                                  <p:cBhvr>
                                    <p:cTn id="8" dur="1" fill="hold">
                                      <p:stCondLst>
                                        <p:cond delay="{dur_ms}"/>
                                      </p:stCondLst>
                                    </p:cTn>
                                    <p:tgtEl>
                                      <p:spTgt spid="{nudge_id}"/>
                                    </p:tgtEl>
                                    <p:attrNameLst>
                                      <p:attrName>style.visibility</p:attrName>
                                    </p:attrNameLst>
                                  </p:cBhvr>
                                  <p:to>
                                    <p:strVal val="visible"/>
                                  </p:to>
                                </p:set>
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0">
                    <p:tgtEl>
                      <p:sldTgt/>
                    </p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>
    """
    slide._element.append(etree.fromstring(xml))


def add_slide_timer(
    slide,
    budget_sec: int,
    elapsed_before_sec: int = 0,
    talk_total_sec: int = 600,
    bar_top=None,
):
    """Add top-right budget badge + countdown bar (auto-starts in Slideshow).

    Returns dict with timing metadata for notes.
    """
    if bar_top is None:
        bar_top = Inches(6.55)

    remaining_after = talk_total_sec - (elapsed_before_sec + budget_sec)
    budget_label = fmt_mmss(budget_sec)
    after_label = fmt_mmss(max(0, remaining_after))
    used_by_end = fmt_mmss(elapsed_before_sec + budget_sec)

    # --- Top-right HUD ---
    hud_left = Inches(10.55)
    hud_top = Inches(0.18)
    hud_w = Inches(2.55)
    hud_h = Inches(0.95)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, hud_left, hud_top, hud_w, hud_h)
    _fill(panel, SURFACE)
    try:
        panel.adjustments[0] = 0.12
    except Exception:
        pass
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, hud_left, hud_top, Inches(0.08), hud_h)
    _fill(accent, GOLD)

    tb = slide.shapes.add_textbox(hud_left + Inches(0.18), hud_top + Inches(0.06), hud_w - Inches(0.28), hud_h - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    _para(tf, "SLIDE BUDGET", size=9, bold=True, color=MUTED, align=PP_ALIGN.LEFT, space_after=0)
    _para(tf, budget_label, size=26, bold=True, color=GOLD, align=PP_ALIGN.LEFT, space_after=0)
    _para(
        tf,
        f"by end {used_by_end} · left {after_label}",
        size=9,
        bold=False,
        color=MUTED,
        align=PP_ALIGN.LEFT,
        space_after=0,
    )

    # --- Countdown track + fill (above progress strip) ---
    track_left = Inches(0.55)
    track_w = Inches(12.2)
    track_h = Inches(0.12)
    track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_left, bar_top, track_w, track_h)
    _fill(track, TRACK)
    try:
        track.adjustments[0] = 0.5
    except Exception:
        pass

    fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, track_left, bar_top, track_w, track_h)
    _fill(fill, GOLD)
    try:
        fill.adjustments[0] = 0.5
    except Exception:
        pass

    # Tiny legend left of bar
    legend = slide.shapes.add_textbox(track_left, bar_top - Inches(0.22), Inches(5.5), Inches(0.22))
    _para(
        legend.text_frame,
        "Countdown starts in Slideshow · red MOVE ON appears when budget ends · does not auto-advance",
        size=8,
        bold=False,
        color=MUTED,
        space_after=0,
    )

    # Over-budget nudge (hidden until countdown completes)
    nudge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.05),
        bar_top - Inches(0.38),
        Inches(2.7),
        Inches(0.32),
    )
    _fill(nudge, DANGER)
    try:
        nudge.adjustments[0] = 0.2
    except Exception:
        pass
    ntf = nudge.text_frame
    ntf.word_wrap = False
    # Clear default empty para then set centered label on the shape itself
    p = ntf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.clear()
    run = p.add_run()
    run.text = "MOVE ON  →"
    _set_run(run, size=12, bold=True, color=TEXT)

    add_countdown_wipe(slide, fill, nudge, max(1000, int(budget_sec * 1000)))

    return {
        "budget_sec": budget_sec,
        "budget_label": budget_label,
        "elapsed_before_sec": elapsed_before_sec,
        "remaining_after_sec": remaining_after,
        "used_by_end_label": used_by_end,
        "after_label": after_label,
    }


def timer_notes_blurb(meta: dict) -> str:
    return (
        f"\n\nON-SLIDE TIMER: Budget {meta['budget_label']} for this slide. "
        f"Gold bar counts down in Slideshow; red MOVE ON appears when budget ends "
        f"(does not auto-advance). Prefer presenter-timer.html on a second screen for a "
        f"live stopwatch that records actual seconds per slide. "
        f"By end of this slide you should be at ~{meta['used_by_end_label']} of 10:00 "
        f"(~{meta['after_label']} remaining)."
    )
