#!/usr/bin/env python3
"""Inject per-slide budget HUDs + countdown bars into an existing PPTX."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from pptx_timer import add_slide_timer, timer_notes_blurb


def apply_timers(path: str, budgets_sec: list[int], talk_total_sec: int = 600) -> None:
    prs = Presentation(path)
    if len(prs.slides) != len(budgets_sec):
        raise ValueError(
            f"Slide count {len(prs.slides)} != budget count {len(budgets_sec)} for {path}"
        )
    elapsed = 0
    for slide, budget in zip(prs.slides, budgets_sec):
        meta = add_slide_timer(
            slide,
            budget_sec=budget,
            elapsed_before_sec=elapsed,
            talk_total_sec=talk_total_sec,
            bar_top=Inches(6.52),
        )
        notes_tf = slide.notes_slide.notes_text_frame
        existing = notes_tf.text or ""
        if "ON-SLIDE TIMER:" not in existing:
            notes_tf.text = existing.rstrip() + timer_notes_blurb(meta)
        elapsed += budget
    prs.save(path)
    print(f"Timers applied → {path} (talk {talk_total_sec // 60}:{talk_total_sec % 60:02d})")


if __name__ == "__main__":
    apply_timers(
        "/workspace/case-study-1-reposition-executive.pptx",
        [75, 80, 80, 75, 85, 75, 65, 65],
    )
    apply_timers(
        "/workspace/case-study-2-ao-offering-executive.pptx",
        [75, 75, 75, 75, 60, 75, 90, 75],
    )
