#!/usr/bin/env python3
"""Short operator KPI summary deck (WhatsApp / interview share)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0F, 0x14, 0x19)
SURFACE = RGBColor(0x1A, 0x23, 0x32)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTED = RGBColor(0x8B, 0x9C, 0xB3)
ACCENT = RGBColor(0x58, 0xA6, 0xFF)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
SUCCESS = RGBColor(0x3F, 0xB9, 0x50)
W = Inches(13.333)
H = Inches(7.5)
OUT = "/workspace/operator-kpi-summary.pptx"


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    sp = bg._element
    s.shapes._spTree.remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def run(r, size=14, bold=False, color=TEXT):
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"


def para(tf, text, size=14, bold=False, color=TEXT, space=6):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text and len(tf.paragraphs) == 1 else tf.add_paragraph()
    if tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1:
        p = tf.paragraphs[0]
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = text
    run(r, size, bold, color)
    return p


def bullets(slide, left, top, width, height, items, size=13):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ""
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "▸  " + item
        run(r, size, False, TEXT)
    return tb


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 Title
    s = blank(prs)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(2.5))
    tf = tb.text_frame
    para(tf, "Telecom Operator KPI Control Tower", 32, True, GOLD, 8)
    para(tf, "Holistic business & operational KPIs · BSS + OSS/IT + Network + Infra/Cloud", 16, False, MUTED, 12)
    para(tf, "Tier-1 operator pattern (illustrative) · Gaurav Chaudhary", 13, False, ACCENT, 0)
    tb2 = s.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(12), Inches(0.5))
    para(tb2.text_frame, "Full catalogue: operator-kpi-catalog.docx  ·  Live sample: operator-kpi-dashboard.html", 12, False, MUTED, 0)

    # 2 Architecture
    s = blank(prs)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
    para(tb.text_frame, "Four-layer KPI architecture", 24, True, GOLD, 0)
    layers = [
        ("1  EXECUTIVE", "12–15 north-stars: revenue, churn, NPS, availability, O2A, bill accuracy, DSO, MTTR"),
        ("2  DOMAIN", "BSS · OSS/IT · Network · Infra/Cloud · Security · Multivendor"),
        ("3  JOURNEY", "Buy → Activate → Use → Bill → Pay → Care → Retain"),
        ("4  DRILL-DOWN", "App, site, vendor, batch, API, cell, cloud account"),
    ]
    for i, (h, b) in enumerate(layers):
        y = Inches(1.5) + i * Inches(1.35)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(12), Inches(1.15))
        box.fill.solid()
        box.fill.fore_color.rgb = SURFACE
        box.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.95), y + Inches(0.15), Inches(11.5), Inches(0.9))
        tf = tb.text_frame
        para(tf, h, 14, True, ACCENT, 2)
        para(tf, b, 13, False, TEXT, 0)

    # 3 BSS examples
    s = blank(prs)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
    para(tb.text_frame, "BSS & revenue-cycle KPIs (your examples + extensions)", 22, True, GOLD, 0)
    bullets(s, Inches(0.7), Inches(1.4), Inches(5.8), Inches(5.5), [
        "Sales: COP · CHQ · digital · partners",
        "Churn: voluntary / involuntary; MNP by operator",
        "Order capture lead time & fallout",
        "Order-to-activate (O2A)",
        "Billing E2E completion; bill prep → confirm",
        "Billing reruns (count + root cause)",
        "Accounts ready for collection",
        "In collection — stage clarity",
        "Pending write-off trigger",
        "DSO · dunning · bad debt",
    ], 13)
    bullets(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5), [
        "Application availability (gold tier)",
        "P1/P2 MTTR · repeat incidents",
        "Change success · emergency rate",
        "API / transaction success (journeys)",
        "Network availability (RAN/core)",
        "Cloud spend vs budget (FinOps)",
        "Patch / backup compliance",
        "Security MTTD / MTTC",
        "Auto-remediation rate",
        "SPOG correlation BSS+OSS+network",
    ], 13)

    # 4 Journey
    s = blank(prs)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
    para(tb.text_frame, "Customer journey → stack mapping", 22, True, GOLD, 0)
    rows = [
        ("BUY", "COP/CHQ sales, conversion", "CRM · channels"),
        ("ACTIVATE", "O2A, port-in, provisioning", "BSS · network"),
        ("USE", "Availability, QoS, complaints", "Apps · RAN · core"),
        ("BILL", "Cycle time, reruns, accuracy", "Billing · mediation"),
        ("PAY", "Collection stages, DSO", "Collections · finance"),
        ("CARE", "FCR, repeat contacts", "CRM · ITSM"),
        ("RETAIN", "Churn by operator, ARPU", "All domains"),
    ]
    y0 = Inches(1.35)
    for i, (j, k, st) in enumerate(rows):
        y = y0 + i * Inches(0.78)
        tb = s.shapes.add_textbox(Inches(0.7), y, Inches(1.5), Inches(0.5))
        para(tb.text_frame, j, 12, True, ACCENT, 0)
        tb = s.shapes.add_textbox(Inches(2.3), y, Inches(5.5), Inches(0.5))
        para(tb.text_frame, k, 12, False, TEXT, 0)
        tb = s.shapes.add_textbox(Inches(8.0), y, Inches(4.5), Inches(0.5))
        para(tb.text_frame, st, 12, False, MUTED, 0)

    # 5 AO maturity
    s = blank(prs)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.6))
    para(tb.text_frame, "Autonomous operations maturity KPIs", 22, True, GOLD, 0)
    bullets(s, Inches(0.7), Inches(1.5), Inches(12), Inches(4), [
        "Reactive → Predictive → Autonomous mix (% tickets / incidents)",
        "% auto-remediation (closed-loop self-heal)",
        "Top-offender closure rate",
        "Human touch rate on P1 scenarios",
        "Playbook / automation reuse across domains",
        "Observability coverage (% critical services with SLO)",
    ], 14)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(12), Inches(1.2))
    tf = tb.text_frame
    para(tf, "End-state visual", 12, True, SUCCESS, 4)
    para(tf, "Open operator-kpi-dashboard.html — sample charts for churn, COP/CHQ sales, billing funnel, availability, journey trends, autonomy path.", 13, False, TEXT, 0)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
